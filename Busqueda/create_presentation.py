from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

def create_tech_presentation():
    # Crear nueva presentación
    prs = Presentation()
    
    # Configurar colores corporativos
    primary_color = RGBColor(46, 125, 50)  # Verde principal
    secondary_color = RGBColor(67, 160, 71)  # Verde secundario
    accent_color = RGBColor(255, 193, 7)  # Amarillo acento
    
    # SLIDE 1: Portada
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    
    # Fondo de color
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = primary_color
    background.line.fill.background()
    
    # Título principal
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "SOLUCIONES TECNOLÓGICAS PARA CHILLÁN"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    p2 = subtitle_frame.paragraphs[0]
    p2.text = "Inteligencia Artificial y IoT para la Región de Ñuble"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    # Información de contacto
    contact_box = slide1.shapes.add_textbox(Inches(2), Inches(6.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.clear()
    p3 = contact_frame.paragraphs[0]
    p3.text = "TechSolutions Chillán\ninfo@techsolutions-chillan.cl\n+56 9 XXXX XXXX"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(255, 255, 255)
    p3.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: Oportunidades del Mercado
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "OPORTUNIDADES DEL MERCADO"
    slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide2.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Estadísticas clave
    p = tf.paragraphs[0]
    p.text = "📊 DATOS CLAVE DE CHILLÁN Y REGIÓN DE ÑUBLE:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    stats = [
        "• 9,007 empresas registradas en la región",
        "• 18.5% del empleo en sector agrícola",
        "• Capital de la agroindustria chilena",
        "• Mercado IA agrícola: $4.7B proyectado para 2028",
        "• Crecimiento del 140% en inversión AgTech (2023-2029)"
    ]
    
    for stat in stats:
        p = tf.add_paragraph()
        p.text = stat
        p.font.size = Pt(16)
        p.level = 1
    
    # Agregar recuadro destacado
    highlight_box = slide2.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    highlight_frame = highlight_box.text_frame
    highlight_frame.clear()
    p_highlight = highlight_frame.paragraphs[0]
    p_highlight.text = "🎯 OPORTUNIDAD: Chillán está posicionada como el hub tecnológico agroindustrial ideal para implementar soluciones de IA avanzadas"
    p_highlight.font.size = Pt(14)
    p_highlight.font.bold = True
    p_highlight.font.color.rgb = RGBColor(255, 255, 255)
    p_highlight.alignment = PP_ALIGN.CENTER
    
    # Fondo del recuadro
    highlight_bg = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(8.4), Inches(1.9)
    )
    highlight_bg.fill.solid()
    highlight_bg.fill.fore_color.rgb = secondary_color
    highlight_bg.line.fill.background()
    # Mover el fondo detrás del texto
    slide2.shapes._spTree.remove(highlight_bg._element)
    slide2.shapes._spTree.insert(-2, highlight_bg._element)
    
    # SLIDE 3: Nuestras Soluciones
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "NUESTRAS SOLUCIONES TECNOLÓGICAS"
    slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide3.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    solutions = [
        {
            "title": "🤖 SISTEMAS RAG PARA AGROINDUSTRIA",
            "description": "Gestión inteligente del conocimiento agrícola con IA generativa"
        },
        {
            "title": "🌤️ MONITOREO ATMOSFÉRICO IoT", 
            "description": "Redes de sensores para calidad del aire y variables ambientales"
        },
        {
            "title": "🎓 IA EN EDUCACIÓN TÉCNICA",
            "description": "Plataformas personalizadas para UBB e INACAP"
        },
        {
            "title": "⚙️ OPTIMIZACIÓN DE PROCESOS",
            "description": "Automatización inteligente para manufactura e industria"
        }
    ]
    
    for i, solution in enumerate(solutions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = solution["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        p_desc = tf.add_paragraph()
        p_desc.text = solution["description"]
        p_desc.font.size = Pt(14)
        p_desc.level = 1
        
        # Espacio entre soluciones
        if i < len(solutions) - 1:
            tf.add_paragraph()
    
    # SLIDE 4: Sectores y Clientes Objetivo
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "SECTORES Y CLIENTES OBJETIVO"
    slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide4.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sectors = [
        {
            "sector": "🚜 AGROINDUSTRIA",
            "companies": ["Vitafoods SpA", "Ideal (Grupo Bimbo)", "Cumfrut", "Colún"]
        },
        {
            "sector": "🏫 EDUCACIÓN",
            "companies": ["Universidad del Bío-Bío", "INACAP Chillán", "Centros de Investigación"]
        },
        {
            "sector": "🏭 MANUFACTURA",
            "companies": ["Parque Industrial Chillán", "Empresas Logísticas", "Sector Maderero"]
        },
        {
            "sector": "🌿 AMBIENTAL",
            "companies": ["SINCA", "Empresas de Monitoreo", "Gestión de Recursos"]
        }
    ]
    
    for i, sector_info in enumerate(sectors):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = sector_info["sector"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        for company in sector_info["companies"]:
            p_company = tf.add_paragraph()
            p_company.text = f"• {company}"
            p_company.font.size = Pt(14)
            p_company.level = 1
        
        # Espacio entre sectores
        if i < len(sectors) - 1:
            tf.add_paragraph()
    
    # SLIDE 5: Beneficios y ROI
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "BENEFICIOS Y RETORNO DE INVERSIÓN"
    slide5.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide5.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Beneficios principales
    p = tf.paragraphs[0]
    p.text = "💰 RETORNO DE INVERSIÓN COMPROBADO:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    benefits = [
        "• Reducción de costos operativos: 20-40%",
        "• Aumento de productividad: 15-35%",
        "• Mejora en calidad de productos: 25%",
        "• ROI promedio: 12-18 meses",
        "• Reducción tiempo de consultas: 40%"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(16)
        p.level = 1
    
    # Separador
    tf.add_paragraph()
    p_cta = tf.add_paragraph()
    p_cta.text = "🚀 PRÓXIMOS PASOS:"
    p_cta.font.size = Pt(18)
    p_cta.font.bold = True
    p_cta.font.color.rgb = primary_color
    
    next_steps = [
        "1. Evaluación gratuita de necesidades tecnológicas",
        "2. Propuesta personalizada con casos de uso específicos",
        "3. Implementación piloto con métricas definidas",
        "4. Escalamiento y optimización continua"
    ]
    
    for step in next_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        p.level = 1
    
    # Agregar call to action final
    cta_box = slide5.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    cta_frame = cta_box.text_frame
    cta_frame.clear()
    p_final = cta_frame.paragraphs[0]
    p_final.text = "📞 ¡Contactanos hoy para transformar tu empresa con IA!"
    p_final.font.size = Pt(18)
    p_final.font.bold = True
    p_final.font.color.rgb = RGBColor(255, 255, 255)
    p_final.alignment = PP_ALIGN.CENTER
    
    # Fondo del CTA
    cta_bg = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(8.4), Inches(1.4)
    )
    cta_bg.fill.solid()
    cta_bg.fill.fore_color.rgb = accent_color
    cta_bg.line.fill.background()
    # Mover el fondo detrás del texto
    slide5.shapes._spTree.remove(cta_bg._element)
    slide5.shapes._spTree.insert(-2, cta_bg._element)
    
    # Guardar presentación
    prs.save('/workspaces/waldo/Presentacion_TechSolutions_Chillan.pptx')
    print("✅ Presentación creada exitosamente: Presentacion_TechSolutions_Chillan.pptx")

if __name__ == "__main__":
    create_tech_presentation()